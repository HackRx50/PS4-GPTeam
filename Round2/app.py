import argparse
import os
from typing import List, Dict
from flask import Flask, request, jsonify
import requests
import google.generativeai as genai
import chromadb
from chromadb.utils import embedding_functions
from sentence_transformers import SentenceTransformer
from tqdm import tqdm
import pdfplumber
from pptx import Presentation
from PIL import Image
import pytesseract
import platform
from flask_cors import CORS 
from transformers import T5ForConditionalGeneration, T5Tokenizer
import warnings
import json
from collections import defaultdict
import cv2
import numpy as np
from pdf2image import convert_from_path
import re
from docx import Document
import random

base_url = "https://hackrx-ps4.vercel.app"
# Initialize conversation state
conversation_state = defaultdict(dict)

# Suppress specific warnings from Hugging Face transformers library
warnings.filterwarnings("ignore", message="You are using the default legacy behaviour")
warnings.filterwarnings("ignore", message="It will be set to `False` by default.")
warnings.filterwarnings("ignore", message="`clean_up_tokenization_spaces` was not set")

# Set Google API key for Gemini
def set_google_api_key():
    api_key = "AIzaSyD7VrRJrSa3W7u0syiZpWldChRCTiWLp-4"
    os.environ["GOOGLE_API_KEY"] = api_key

set_google_api_key()

# Define the default actions list outside the function
DEFAULT_ACTIONS_LIST = ["create_order", "cancel_order", "collect_payment", "view_invoice"]

# Initialize Flask app
app = Flask(__name__)
CORS(app)
collection = None

order_id = 12345  # Static order ID
mobile = "555-1234"  # Static mobile number
action = "create_order"  # Static action

# Now you can retrieve the API key from the environment variable when needed
google_api_key = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=google_api_key)

UPLOAD_FOLDER = os.path.join(os.getcwd(), 'uploads')
if not os.path.exists(UPLOAD_FOLDER):
    print("Creating folder...")
    os.makedirs(UPLOAD_FOLDER)
else:
    print("Folder already exists.")

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Load the Flan-T5 model and tokenizer from Hugging Face
flan_tokenizer = T5Tokenizer.from_pretrained("google/flan-t5-large", legacy=False)
flan_model = T5ForConditionalGeneration.from_pretrained("google/flan-t5-large")

# Initialize history storage
session_history: Dict[str, List[Dict[str, str]]] = {}

# Load Huggingface model for embeddings
embedding_model = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")

# Initialize Gemini model
model = genai.GenerativeModel("gemini-1.5-flash")

# Helper functions to extract text from different document types


def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file with layout preservation and OCR fallback for scanned PDFs."""
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text(layout=True)  # Preserve layout
            if page_text:  
                text += page_text + "\n"
            else:  # If text extraction fails, use OCR on the image version of the page
                images = convert_from_path(pdf_path)
                for image in images:
                    ocr_text = pytesseract.image_to_string(image)
                    text += ocr_text + "\n"
    print("EXTRACTED TEXT: \n", clean_text(text))
    return clean_text(text)

def extract_text_from_ppt(ppt_path):
    """Extract text from a PPT file, including text in grouped shapes and tables."""
    prs = Presentation(ppt_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
            # Extract text from tables
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " ".join([cell.text for cell in row.cells])
                    text += row_text + "\n"
            # Check for grouped shapes
            if hasattr(shape, "shapes"):
                for sub_shape in shape.shapes:
                    if hasattr(sub_shape, "text"):
                        text += sub_shape.text + "\n"
    print("EXTRACTED TEXT: \n", clean_text(text))
    return clean_text(text)

def extract_text_from_docx(file_path):
    """Extract text from a .docx file."""
    doc = Document(file_path)
    text = '\n'.join([para.text for para in doc.paragraphs])
    print("EXTRACTED TEXT: \n", clean_text(text))
    return clean_text(text)

def preprocess_image(image):
    """Preprocess the image for better OCR accuracy: grayscale, contrast, and resize."""
    gray = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2GRAY)
    # Resize to improve OCR accuracy if image is too small
    resized = cv2.resize(gray, None, fx=2, fy=2, interpolation=cv2.INTER_CUBIC)
    # Increase contrast
    contrast = cv2.convertScaleAbs(resized, alpha=2.0, beta=0)
    return contrast

def extract_text_from_image(image_path):
    """Extract text from an image using OCR with preprocessing."""
    image = Image.open(image_path)
    preprocessed_image = preprocess_image(image)  # Apply preprocessing for better accuracy
    text = pytesseract.image_to_string(preprocessed_image, config='--psm 6')  # Use custom OCR config for block text
    print("EXTRACTED TEXT: \n", clean_text(text))
    return clean_text(text)

def extract_text_from_txt(file_path):
    """Extract text from a text file."""
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def extract_text_from_file(file_path):
    """Extract text based on the file type."""
    _, extension = os.path.splitext(file_path)
    
    if extension.lower() == ".pdf":
        return extract_text_from_pdf(file_path)
    elif extension.lower() in [".ppt", ".pptx"]:
        return extract_text_from_ppt(file_path)
    elif extension.lower() in [".jpg", ".jpeg", ".png", ".bmp", ".tiff"]:
        return extract_text_from_image(file_path)
    elif extension.lower() == ".txt":
        return extract_text_from_txt(file_path)
    elif extension.lower() == ".docx":
        return extract_text_from_docx(file_path)
    elif extension.lower() == ".json":
        # Open the file and load the JSON data with UTF-8 encoding
        extracted_text = extract_text_JSON_from_file(file_path)
        return clean_text(extracted_text)
    else:
        return None


def extract_text_JSON(json_object, indent_level=0):
    """Extract text from a JSON object with indentation."""
    lines = []
    indent = '    ' * indent_level  # Create indentation for better readability

    if isinstance(json_object, dict):
        for key, value in json_object.items():
            if isinstance(value, (dict, list)):
                # For nested dictionaries or lists, call the function recursively
                nested_value = extract_text_JSON(value, indent_level + 1)
                lines.append(f'{indent}{key}: {nested_value.strip()}')  # Append key with nested values
            else:
                # Directly append values, with indentation
                lines.append(f'{indent}{key}: {value}')  # Append the key-value pair directly
    elif isinstance(json_object, list):
        for index, item in enumerate(json_object):
            if isinstance(item, (dict, list)):
                # For nested dictionaries or lists, call the function recursively
                nested_value = extract_text_JSON(item, indent_level + 1)
                lines.append(f'{indent}Item {index + 1}: {nested_value.strip()}')  # Append with item index
            else:
                # Directly append values
                lines.append(f'{indent}Item {index + 1}: {item}')  # Append the item directly

    # Join lines into a single string with a period after each line
    return '\n'.join(lines) + '.'  # Ensure there's a period at the end of the string

def extract_text_JSON_from_file(file_path):
    """Extract text from a JSON file."""
    with open(file_path, 'r', encoding='utf-8') as json_file:
        # Load the entire JSON file
        json_object = json.load(json_file)  # This will load the whole JSON object into memory
        return extract_text_JSON(json_object)  # Call the extraction function


def clean_text(text):
    """Remove unnecessary empty lines and extra spaces from extracted text."""
    # Remove leading/trailing spaces and collapse multiple spaces into one
    cleaned_text = re.sub(r'\s+', ' ', text.strip())
    # Remove multiple empty lines and keep only one
    cleaned_text = re.sub(r'\n\s*\n', '\n', cleaned_text)
    return cleaned_text


def search_json_for_keys(data, keys):
    """
    Recursively search for the given keys in a nested JSON structure.
    :param data: The JSON data (could be dict or list)
    :param keys: The list of keys to search for
    :return: A dictionary with found key-value pairs
    """
    found = {}
    if isinstance(data, dict):
        for key, value in data.items():
            if key in keys:
                found[key] = value  # Found the key, add to result
            elif isinstance(value, (dict, list)):
                # Recursively search in nested structures
                found.update(search_json_for_keys(value, keys))
    elif isinstance(data, list):
        for item in data:
            found.update(search_json_for_keys(item, keys))
    
    return found


def execute_action(action_name: str, query: str, session_id: str, document_id: str = None) -> str:
    """
    Executes the specified action based on the provided action name and query.
    Dynamically retrieves missing information from a collection if needed.
    """
    print("action", action_name)
    print("query", query)
    # Extract key points using Gemini
    key_points_json = extract_key_action_with_gemini(query)
    extracted_data = json.loads(key_points_json)
    print("Extracted data:", extracted_data)

    if action_name == "eligibility_check":
        response = eligibility_check(mobile)
        return f"The eligibility check for mobile {mobile} has been completed. Results: {response}"
    elif action_name == "health_check":
        response = health_check()
        return f"The health check has been conducted. All systems are operational. Response: {response}"
    elif action_name == "generate_lead":
        response = generate_lead(mobile)
        return f"A lead has been successfully generated for mobile {mobile}. Response: {response}"
    
    # Extracted generic data from the query
    entity_id = extracted_data.get("id")  # Generic id
    name = extracted_data.get("name")  # Generic name
    amount = extracted_data.get("amount")  # Generic amount
    price = extracted_data.get("price")  # Generic price

    print(f"id: {entity_id}, Name: {name}, Amount: {amount}, Price: {price}")

    # Check for missing information based on the action type
    missing_info = []
    if action_name == "create_entity" :
        if entity_id is None:
            missing_info.append("id")
        if name is None:
            missing_info.append("name")
        missing_info.append("id")
    if action_name == "cancel_entity" and entity_id is None:
        missing_info.append("id")
    if action_name == "process_payment" and amount is None:
        missing_info.append("amount")
    if action_name == "view_entity" and entity_id is None:
        missing_info.append("id")

    print("Missing info:", missing_info)
    missing_items_string = ", ".join(missing_info)
    # If there is any missing info, formulate a dynamic query to retrieve it
    if missing_info:
        for missing_field in missing_info:
            missing_query = f"What is the {missing_field.lower()}?"

            # Query the collection to retrieve the missing info
            context_results = collection.query(
                query_texts=[missing_query],   # Dynamic query
                n_results=7,
                include=["documents", "metadatas"],
                where={"document_id": document_id}  # Filter based on document id
            )
            print("Context results:", context_results)

            # Process query results and extract the missing information
            for result in context_results.get('documents', []):
                if isinstance(result, list) and len(result) > 0:
                    document_text = result[0]  # Access the string from the list
                    print("Document Text:", document_text)
                    # Use Gemini to extract relevant information from the document
                    gemini_response = extract_key_action_with_gemini_missing(document_text,missing_info)
                    extracted_data = json.loads(gemini_response)
                    print("Extracted Data from Gemini:", extracted_data)

                    # Search for the required missing information using the recursive function
                    found_data = search_json_for_keys(extracted_data, missing_info)
                    print("Found Data:", found_data)

                    # Update the missing information if found
                    if 'id' in missing_info and 'id' in found_data:
                        entity_id = found_data['id']
                        missing_info.remove('id')

                    if 'amount' in missing_info and 'amount' in found_data:
                        amount = found_data['amount']
                        missing_info.remove('amount')

                    if 'name' in missing_info and 'name' in found_data:
                        name = found_data['name']
                        missing_info.remove('name')
                else:
                    print(f"Unexpected data format in query result: {result}")

        # If missing info is still not retrieved, prompt the user for manual input
        if missing_info:
            conversation_state[session_id]["pending_action"] = action_name
            conversation_state[session_id]["missing_info"] = missing_info
            conversation_state[session_id]["query"] = query
            conversation_state[session_id]["document_id"] = document_id
            conversation_state[session_id]["order_id"] = entity_id if action_name in ["create_order", "cancel_order"] else None
            conversation_state[session_id]["amount"] = amount if action_name == "process_payment" else None
            return f"Please provide the following information: {', '.join(missing_info)}."

    # Proceed with the action if all required information is present
    confirmation_message = ""
    if action_name == "create_order":
        if random2 is None:
            random2 = random.randint(1000, 9999)
        confirmation_message = create_order(order_id=random2,product_id=entity_id, mobile=mobile, product_name=name, product_price=100, action="create_order")
        #confirmation_message = f"Entity with id {entity_id}, name {name}, and price {price} has been created."
    elif action_name == "cancel_order":
        confirmation_message = f"Entity with id {entity_id} has been cancelled."
    elif action_name == "process_payment":
        confirmation_message = f"Payment of amount {amount} has been processed."
    elif action_name == "view_entity":
        confirmation_message = f"Here is the entity with id {entity_id}."
    elif action_name == "eligibility_check":
        response = eligibility_check(mobile)
        return f"The eligibility check for mobile {mobile} has been completed. Results: {response}"
    elif action_name == "health_check":
        response = health_check()
        return f"The health check has been conducted. All systems are operational. Response: {response}"
    elif action_name == "generate_lead":
        response = generate_lead(mobile)
        return f"A lead has been successfully generated for mobile {mobile}. Response: {response}"
    elif action_name == "get_order_status":
        response = get_order_status(entity_id, mobile)
        return f"The order status for order {entity_id} has been retrieved. Response: {response}"
    elif action_name == "get_orders":
        response = get_order_status(entity_id, mobile)
        return f"The order status for order {entity_id} has been retrieved. Response: {response}"
    else:
        confirmation_message = "No action taken."
    return confirmation_message
def confirm_action(action_name: str, identifier: str, mobile: str = None, product_id: str = None, product_name: str = None, product_price: float = None) -> str:
    """
    Confirms the specified action and executes it if confirmed.
    Provides detailed confirmation messages based on the action type.
    """

    if action_name == "cancel_order":
        # You can implement the cancel_order function as needed.
        return f"Your order with ID {identifier} has been successfully cancelled."
    
    elif action_name == "create_order":
        response = create_order(identifier, mobile, product_id, product_name, product_price)
        return f"Your order with ID {identifier} has been created successfully. Response: {response}"

    elif action_name == "collect_payment":
        # Implement the collect_payment function as needed and call it here.
        return f"Payment with ID {identifier} has been collected successfully."

    elif action_name == "view_invoice":
        # Implement the view_invoice function as needed and call it here.
        return f"Here is your invoice with ID {identifier}. Please review it carefully."

    elif action_name == "eligibility_check":
        response = eligibility_check(mobile)
        return f"The eligibility check for mobile {mobile} has been completed. Results: {response}"

    elif action_name == "health_check":
        response = health_check()
        return f"The health check has been conducted. All systems are operational. Response: {response}"

    elif action_name == "generate_lead":
        response = generate_lead(mobile)
        return f"A lead has been successfully generated for mobile {mobile}. Response: {response}"

    elif action_name == "get_orders":
        response = get_order_status(identifier, mobile)
        return f"Here are the orders associated with ID {identifier}. Response: {response}"

    elif action_name == "get_order_status":
        response = get_order_status(identifier, mobile)
        return f"The status for the order with ID {identifier} has been retrieved. Response: {response}"

    # If the action name does not match any known category
    return "No action taken."

# Define the default actions list outside the function
DEFAULT_ACTIONS_LIST = ["create_order", "cancel_order", "collect_payment", "view_invoice"]

def classify_query_with_flan(query: str, actions_list: list = DEFAULT_ACTIONS_LIST) -> str:
    """Use Flan-T5 to classify the query as action-based or context-based using a dynamic list of actions."""
    
    # Convert the list of actions to a string to use in the prompt
    actions_str = ', '.join(actions_list)
    
    # Dynamic prompt using the actions list
    prompt = f"""
    Classify the following query strictly as one of the actions: {actions_str}, or context-based.
    
    Query: {query}
    """
    
    # Tokenize the input
    inputs = flan_tokenizer(prompt, return_tensors="pt")
    outputs = flan_model.generate(**inputs, max_length=10, num_return_sequences=1)
    response = flan_tokenizer.decode(outputs[0], skip_special_tokens=True).lower()
    print(response)
    
    # Check if the response matches one of the actions
    for action in actions_list:
        if action in response:
            return action
    return "context_based"


def build_combined_prompt(query: str, context: List[str], history: List[Dict[str, str]]) -> str:
    base_prompt = """I am going to ask you a question, and your answer should be based strictly on the context provided from the document, along with the history of our previous interactions. 
Please extract all relevant information from the document to generate a comprehensive and accurate response, while also taking into account the knowledge from the responses you have already provided during this conversation.
If the context is insufficient or unclear, make a reasonable inference and a best guess based on both the document content and the history of our conversation. Ensure that your response is directly aligned with the user query and that you use both the document as the source of truth and the stored conversation history.
Your response must be informative, provide as much relevant detail as possible, and never leave the question unanswered unless absolutely necessary. If any additional key information is missing, prompt the user for further clarification.
        """
    user_prompt = f"The question is '{query}'. Here is all the context you have: {' '.join(context)}"
    #history_prompt = "\n".join([f"User: {item['query']}\nBot: {item['response']}" for item in history])
    history_prompt= "Here is all the history you have:\n" + "\n".join([f"User: {item['query']}\nBot: {item['response']}" for item in history])
    history_prompt = clean_text(history_prompt)
    print()
    print("HISTORY : ", f"{history_prompt}")

    return f"{base_prompt} {history_prompt} {user_prompt}"

# Common headers
headers = {
    "x-team": "GPTeam"
}

# 1. Generate Lead
def generate_lead(mobile):
    url = f"{base_url}/generate-lead"
    data = {"mobile": mobile}
    response = requests.post(url, json=data, headers=headers)
    return response.json()

# 2. Eligibility Check
def eligibility_check(mobile):
    url = f"{base_url}/eligibility-check"
    data = {"mobile": mobile}
    response = requests.post(url, json=data, headers=headers)
    return response.json()

# 3. Health Check
def health_check():
    url = f"{base_url}/health-check"
    response = requests.get(url, headers=headers)
    return response.json()

# Static values for order_id, mobile, and action


# 4. Create Order
def create_order(order_id, mobile, product_id, product_name, product_price, action="create_order"):
    url = f"{base_url}/order"
    data = {
        "id": order_id,
        "mobile": mobile,
        "productId": product_id,
        "productName": product_name,
        "productPrice": product_price,
        "action": action  # Include action parameter
    }
    response = requests.post(url, json=data, headers=headers)
    print(response.json())
    return response.json()

# 5. Get All Orders
def get_orders():
    url = f"{base_url}/orders"
    response = requests.get(url, headers=headers)
    return response.json()

# 6. Get Order Status
def get_order_status(order_id, mobile):
    url = f"{base_url}/order-status"
    params = {"orderId": order_id, "mobile": mobile}
    response = requests.get(url, params=params, headers=headers)
    return response.json()

def create_questions(query: str) -> str:
    # Create a nuanced and engaging prompt for the model
    prompt = f"""
    You are a sophisticated AI assistant trained to convert statements into concise questions. Your task is to analyze the following statement and determine the best possible question that reflects the user's intent.

    **Instructions**:
    1. If the input is already a question, return it unchanged.
    2. If the input is a statement, transform it into a clear and concise question that captures the essence of what the user might want to know.
    3. If the statement is unclear or does not indicate a specific user intent, generate a question that reflects a possible inquiry based on the given information. Make it the best possible question from the query.
    4. Ensure that your question is grammatically correct and contextually appropriate.
    5. Do not include any explanations or additional commentary—just provide the question.

    **Statement**: {query}

    For example: 
    - If the input is "Document ID for the report," the output should be "What is the document ID for the report?"
    - If the input is "title," the output should be "What is the title of the document?"
    """
    
    # Generate the response using the model
    response = model.generate_content(prompt)
    
    # Extract and clean the response
    classification = response.text.strip()
    
    return classification

def extract_key_action_with_gemini(query: str) -> str:
    """
    Use Gemini to extract key information from the query.
    """
    prompt = f"""
    You are a helpful assistant. Extract key information from the statement below, mapping various types of data into generic parameters. Use the following mapping guidelines:

    - Any type of identifier (e.g., order ID, invoice ID, product ID, customer ID) should be mapped to 'ID'.
    - Any type of name (e.g., product name, customer name, entity name) should be mapped to 'name'.
    - Any type of monetary amount (e.g., total, balance, payment) should be mapped to 'amount'.
    - Any price-related information should be mapped to 'price'.
    - Any additional key-value pairs should be extracted and assigned generic parameter names based on their context.

    Please return the extracted information in a well-structured JSON format.

    Statement: "{query}"
    """

    # Send the prompt to Gemini for classification
    response = model.generate_content(prompt)
    classification = response.text.strip().lower()
    #print("Gemini Key information before cleaning:", classification)
    classification = classification.replace("json", "").strip() # Remove "json" from the string and strip leading/trailing whitespace
    while classification.startswith('"') or classification.startswith("'") or classification.startswith("`"):
        classification = classification[1:]  # Remove the first character if it's a quote
    while classification.endswith('"') or classification.endswith("'") or classification.endswith("`"):
        classification = classification[:-1]  # Remove the last character if it's a quote

    # Remove backticks from the string
    classification = classification.replace("`", "")
    print("Gemini Key information after cleaning:", classification)
     # Convert string to dictionary to ensure valid JSON and remove any null values
    try:
        extracted_data = json.loads(classification)
        filtered_data = {k: v for k, v in extracted_data.items() if v is not None}
        return json.dumps(filtered_data)  # Return as JSON without null values
    except json.JSONDecodeError:
        print("JSON contains null values")
        return "{}"
def extract_key_action_with_gemini_missing(query: str, missing_items_string: list) -> str:
    """
    Use Gemini to extract key information from the query.
    """
    prompt = f"""
    You are a helpful assistant. Extract key informations {missing_items_string} from the statement below, with missing items being the key and value being extracted info. 

    - Map any type of identifier (e.g., order ID, invoice ID, product ID, customer ID) to 'ID'.
    - Map any type of name (e.g., product name, customer name, entity name) to 'name'.
    - Map any type of monetary amount (e.g., total, balance, payment) to 'amount'.
    - Map any price-related information to 'price'.
    - Additionally, check the following list for expected items that should be mentioned in the statement: {missing_items_string}.
    - Extract which items from this list are missing in the statement and include them in a 'missing' list in the output.

    Return the extracted information in a well-structured JSON format. REMEMBER YOU SHOULD EXTRACT ONLY {missing_items_string} AND NOT ANYTHING THAN THAT

    Statement: "{query}"
"""

    # Send the prompt to Gemini for classification
    response = model.generate_content(prompt)
    classification = response.text.strip().lower()
    #print("Gemini Key information before cleaning:", classification)
    classification = classification.replace("json", "").strip() # Remove "json" from the string and strip leading/trailing whitespace
    while classification.startswith('"') or classification.startswith("'") or classification.startswith("`"):
        classification = classification[1:]  # Remove the first character if it's a quote
    while classification.endswith('"') or classification.endswith("'") or classification.endswith("`"):
        classification = classification[:-1]  # Remove the last character if it's a quote

    # Remove backticks from the string
    classification = classification.replace("`", "")
    print("Gemini Key information after cleaning:", classification)
     # Convert string to dictionary to ensure valid JSON and remove any null values
    try:
        extracted_data = json.loads(classification)
        filtered_data = {k: v for k, v in extracted_data.items() if v is not None}
        return json.dumps(filtered_data)  # Return as JSON without null values
    except json.JSONDecodeError:
        print("JSON contains null values")
        return "{}"
    #return classification
def classify_query_with_gemini(query: str) -> str:
    order_functions = [
        "create_order",
        "cancel_order",
        "collect_payment",
        "view_invoice",
        "eligibility_check",
        "health_check",
        "generate_lead",
        "get_orders",
        "get_order_status",
        "context_based",
        "context-based"
    ]

    prompt = f"""
    You are a Classification expert. Strictly Classify the following query into one of these categories:
1. create_order
2. cancel_order
3. collect_payment
4. view_invoice
5. eligibility_check
6. health_check
7. generate_lead
8. get_orders
9. get_order_status
10. context_based

Rules for classification:
- If the query explicitly mentions creating an order, classify it as 'create_order'.
- If the query explicitly mentions cancelling an order, classify it as 'cancel_order'.
- If the query is about making a payment or collecting money, classify it as 'collect_payment'.
- If the query is about viewing or requesting an invoice, classify it as 'view_invoice'.
- If the query pertains to checking eligibility for a service or program or anything else, classify it as 'eligibility_check'.
- If the query is about checking the status of an application or system, classify it as 'health_check'.
- If the query seeks to generate a lead for potential sales or contacts, classify it as 'generate_lead'.
- If the query is about retrieving all orders and **explicitly mentions an order ID**, classify it as 'get_orders'.
- If the query asks for the status of a specific order, classify it as 'get_order_status'.
- If no specific order ID is mentioned for retrieving orders or statuses, classify the query as 'context_based'.
- If the query is not of any of the first 9 categories, classify it as context_based.

Examples:
- "I want to place an order" -> create_order
- "Cancel my recent order" -> cancel_order
- "How do I pay for my order?" -> collect_payment
- "Can I see my invoice?" -> view_invoice
- "Am I eligible for this promotion?" -> eligibility_check
- "Is the server running smoothly?" -> health_check
- "I need more information on your services" -> generate_lead
- "Show me all my orders with ID #12345" -> get_orders
- "What is the status of my order #12345?" -> get_order_status
- "Show me my last order" -> context_based (since no order ID is mentioned)
- "What's your return policy?" -> context_based
- "What is my last order?" -> context_based

Query: "{query}"

Classification:
    """

    # Send the prompt to Gemini for classification
    response = model.generate_content(prompt)
    classification = response.text.strip().lower()
    print(classification)
    for i in order_functions:  
        if i in classification:
            print("Gemini Classification:", i)
            if i == "context-based":
                return "context_based"
            return i
def prompt_gemini(query: str) -> str:
    # Send the prompt to Gemini for classification
    response = model.generate_content(query)
    classification = response.text.strip().lower()
    print(classification)
    return classification
def get_gemini_response(query: str, context: List[str], session_id: str, extract_data: bool = False,documentID: str = None) -> str:
    history = session_history.get(session_id, [])

    # Classify the query using Flan-T5
    #action = classify_query_with_flan(query)
    action = classify_query_with_gemini(query)
    action = action.lower()
    if action != "context_based":
        # If action is identified, extract necessary details
        action_response = execute_action(action, query,session_id, document_id=documentID)
        # Append the action execution to session history
        session_history.setdefault(session_id, []).append({"query": query, "response": action_response})
        return True, action_response
    print("\n\n\n ")
    print(history)
    # Build the combined prompt
    prompt = build_combined_prompt(query, context, history)
    #print("SENT PROMPT TO GEMINI :", prompt)
    # Get response from Gemini
    response = model.generate_content(prompt)
    response_text = response.text.strip().lower()
    # Save the query and response in session history
    session_history.setdefault(session_id, []).append({"query": query, "response": response.text})
    return False, response.text


@app.route('/upload_document', methods=['POST'])
def upload_document():
    if 'document' not in request.files:
        return jsonify({"error": "No document uploaded."}), 400

    document = request.files['document']
    doc_name = document.filename

    # Validate document format
    if not doc_name.endswith(('.pdf', '.ppt', '.docx', '.jpg', '.jpeg', '.png', '.txt', '.json', '.csv','.txt','.pptx')):
        return jsonify({"error": "Invalid document format. Only PDF, PPT, DOCX, and images are supported."}), 400
    
        #print(data)
    try:
        # Save the uploaded file temporarily
        upload_path = os.path.join("uploads", doc_name)
        document.save(upload_path)
        

        # Extract text from the uploaded document
        text = extract_text_from_file(upload_path)
        print(text)
        if text is None:
            return jsonify({"error": f"Unsupported document format or failed to extract text from {doc_name}."}), 400
        elif text == "JSON_DATA":
            return jsonify({"Success": f"Action taken successfully.."}), 200
        # Generate a unique document ID
        document_id = "doc_" + os.urandom(6).hex()

        # Split text into lines for processing
        lines = text.splitlines()
        print("SEPARATED: ")
        for line in lines:
            print(line)

        # Store extracted text into ChromaDB with metadata
        documents = [line for line in lines if line.strip()]
        #metadatas = [{"filename": doc_name, "line_number": i} for i, _ in enumerate(documents, 1)]
        metadatas = [{"document_id": document_id, "filename": doc_name, "line_number": i} for i, _ in enumerate(documents, 1)]

        # Add documents to Chroma collection
        collection.add(
            ids=[f"{document_id}_{i}" for i in range(len(documents))],
            documents=documents,
            metadatas=metadatas
        )

        return jsonify({"message": "Document uploaded and processed successfully.", "document_id": document_id}), 200

    except Exception as e:
        return jsonify({"error": f"An error occurred: {str(e)}"}), 500

@app.route('/chat', methods=['POST'])
def chat():
    data = request.json
    user_id = request.headers.get('x-user-id')
    session_id = request.headers.get('x-session-id')
    query = data.get('query')
    document_id = data.get('document_id')  # Optional document ID
    if not query or not isinstance(query, str):
        return jsonify({"error": "Query is required and must be a string oooh."}), 400
    if query is None:
        return jsonify({"error": "Query is missing or None."}), 400
    # Validate that the query is provided and is a string
    if not query or not isinstance(query, str):
        return jsonify({"error": "Query is required and must be a string."}), 400
    print("QUESTIONS CREATED:")
    print(create_questions(query))
    query = create_questions(query)
    try:
        # Handle pending actions
        if session_id in conversation_state and "pending_action" in conversation_state[session_id]:
            pending_action = conversation_state[session_id]["pending_action"]
            order_id = conversation_state[session_id].get("order_id")
            payment_id = conversation_state[session_id].get("payment_id")
            invoice_id = conversation_state[session_id].get("invoice_id")
            print("query:" , query)
            print(type(query))
            if query.lower() in ["yes", "confirm"]:
                response = confirm_action(pending_action, order_id if pending_action in ["cancel_order", "create_order"] else payment_id)
                # Clear pending action after confirming
                del conversation_state[session_id]
                return jsonify({"bot_message": response}), 200

            elif query.lower() in ["no", "cancel"]:
                response = "Action cancelled."
                # Clear pending action
                del conversation_state[session_id]
                return jsonify({"bot_message": response}), 200

        # Fetch context based on document_id
        context_results = collection.query(
            query_texts=[query],
            n_results=5,
            include=["documents", "metadatas"],
            where={"document_id": document_id} if document_id else None  # Use None if no document_id
        )

        # If no relevant documents are found
        if not context_results["documents"]:
            return jsonify({"error": "No relevant information found in the database."}), 404

        # Flattening documents and metadata
        context = [doc for docs in context_results["documents"] for doc in docs]  # Flatten documents
        metadata = [meta for metas in context_results["metadatas"] for meta in metas]  # Flatten metadata

        # Check if context is a flat list of strings
        if not all(isinstance(c, str) for c in context):
            return jsonify({"error": "Context is not in the expected format."}), 500

        # Get response from Gemini model
        action_flag, response = get_gemini_response(query, context, session_id, documentID=document_id)

        if action_flag:
            return jsonify({"bot_message": response}), 200

        # Prepare reference excerpts and metadata
        references = [{
            "excerpt": doc[:200],  # Show the first 200 characters of the relevant context
            "line_number": meta["line_number"],
            "document_name": meta["filename"]
        } for doc, meta in zip(context, metadata)]

        # Format the response with relevant metadata
        response_data = {
            "bot_message": response,
            "references": references  # Include reference excerpts, line numbers, and document names
        }

        return jsonify(response_data), 200

    except Exception as e:
        return jsonify({"error": f"An internal error occurred: {str(e)}"}), 500


def main(collection_name: str = "documents_collection", persist_directory: str = ".") -> None:
    global collection  # Declare the global variable first

    client = chromadb.PersistentClient(path=persist_directory)

    # Create embedding function using Huggingface transformers
    embedding_function = embedding_functions.HuggingFaceEmbeddingFunction(
        model_name="sentence-transformers/all-MiniLM-L6-v2",
        api_key="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"
    )

    try:
        collection = client.get_collection(name=collection_name, embedding_function=embedding_function)
        print(f"Loaded existing collection: {collection_name}")
    except Exception as e:
        print(f"Collection '{collection_name}' not found. Creating a new collection.")
        collection = client.create_collection(name=collection_name, embedding_function=embedding_function)

    # Start the Flask app
    app.run(port=3000)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Load documents into a Chroma collection")
    parser.add_argument("--persist_directory", type=str, default="chroma_storage", help="Directory to store the Chroma collection")
    parser.add_argument("--collection_name", type=str, default="documents_collection", help="Name of the Chroma collection")

    args = parser.parse_args()

    main(collection_name=args.collection_name, persist_directory=args.persist_directory)
